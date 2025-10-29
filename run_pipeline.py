"""
run_pipeline.py

Usage:
    python run_pipeline.py --input sample.pdf --n_slides 8

What it does (best-effort, offline):
1. Extracts text from PDF (pdfminer.six).
2. Splits text into paragraphs and heuristically finds headings.
3. Chooses top-N sections using either sentence-transformers (if installed) + KMeans or Sumy fallback.
4. Produces short headlines (6-20 words), 1-2 bullets, and one-sentence speaker notes per slide.
5. Generates simple illustrative images programmatically (PIL).
6. Assembles a PPTX using python-pptx in ./output/<pdfname>.pptx
7. Generates TTS (pyttsx3) and background music (synthesized), makes MP4 using moviepy with simple Ken Burns effect.

Requirements (pip):
    pip install pdfminer.six python-pptx pillow moviepy pyttsx3 numpy scipy scikit-learn sentence-transformers nltk sumy

Notes:
- This script avoids web downloads and uses programmatic visuals.
- If offline TTS fails, slides will still be turned into a silent video with on-screen captions.

"""

import argparse
import os
import sys
import math
import tempfile
from pathlib import Path
from io import BytesIO

# PDF text extraction
from pdfminer.high_level import extract_text

# PPTX
from pptx import Presentation
from pptx.util import Inches, Pt

# Images
from PIL import Image, ImageDraw, ImageFont

# Audio / Video
from moviepy import ImageClip, AudioFileClip, CompositeAudioClip, concatenate_videoclips

# TTS
try:
    import pyttsx3
    TTS_AVAILABLE = True
except Exception:
    TTS_AVAILABLE = False

# Clustering / embeddings
USE_SBERT = False
try:
    from sklearn.cluster import KMeans
    from sentence_transformers import SentenceTransformer
    USE_SBERT = True
except Exception:
    USE_SBERT = False

# Sumy fallback
SUMY_AVAILABLE = False
try:
    from sumy.parsers.plaintext import PlaintextParser
    from sumy.nlp.tokenizers import Tokenizer
    from sumy.summarizers.lex_rank import LexRankSummarizer
    SUMY_AVAILABLE = True
except Exception:
    SUMY_AVAILABLE = False

# Utilities
import re
import numpy as np

def ensure_outdir(base_out):
    os.makedirs(base_out, exist_ok=True)


def extract_paragraphs(pdf_path):
    text = extract_text(pdf_path)
    # Normalize line endings
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    # Split by two or more newlines as paragraph boundaries
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    return paras


def heuristics_find_headings(paras):
    # Return indices of paragraphs that look like headings
    headings = []
    for i, p in enumerate(paras):
        words = p.split()
        if len(words) <= 8 and (p.isupper() or p.istitle() or re.match(r'^[A-Z][^\.\n]{2,}$', p)):
            # short paragraph & looks like title-case or uppercase
            headings.append(i)
    return headings


def top_n_sections_with_embeddings(paras, n):
    # Use SBERT to embed paragraphs and KMeans to cluster into n clusters.
    model = SentenceTransformer('all-MiniLM-L6-v2')
    embeddings = model.encode(paras, show_progress_bar=False)
    k = min(n, len(paras))
    kmeans = KMeans(n_clusters=k, random_state=42)
    labels = kmeans.fit_predict(embeddings)
    centers = kmeans.cluster_centers_
    chosen = []
    for ci in range(k):
        idxs = [i for i, lab in enumerate(labels) if lab == ci]
        # pick paragraph closest to center
        best = min(idxs, key=lambda i: np.linalg.norm(embeddings[i] - centers[ci]))
        chosen.append((best, paras[best]))
    # sort by original order
    chosen.sort(key=lambda x: x[0])
    return [text for (_, text) in chosen]


def top_n_sections_with_sumy(paras, n):
    text = '\n\n'.join(paras)
    parser = PlaintextParser.from_string(text, Tokenizer('english'))
    summarizer = LexRankSummarizer()
    # LexRank returns sentences; we will join top n sentences
    summary = summarizer(parser.document, n)
    return [str(s) for s in summary]


def make_headline_and_bullets(section_text):
    # Simple heuristic: split into sentences and pick the most informative sentence as headline
    sentences = re.split(r'(?<=[.!?])\s+', section_text.strip())
    sentences = [s.strip() for s in sentences if s.strip()]
    if not sentences:
        return ("", [])
    # Choose the longest sentence as the headline candidate but trim to 6-20 words
    headline = max(sentences, key=lambda s: len(s.split()))
    words = headline.split()
    if len(words) > 20:
        headline = ' '.join(words[:18]) + '...'
    # Make bullets: next 1-2 sentences or key phrases
    bullets = []
    for s in sentences[:3]:
        if s != headline:
            bullets.append(s if len(s.split()) <= 20 else ' '.join(s.split()[:20]) + '...')
        if len(bullets) >= 2:
            break
    if not bullets and len(sentences) > 1:
        bullets = [sentences[1]]
    # Speaker note: one sentence (first non-headline)
    speaker = bullets[0] if bullets else sentences[0]
    return (headline, bullets, speaker)


from PIL import Image, ImageDraw, ImageFont, ImageFilter
import textwrap
import os

def generate_simple_image(text, out_path, bullets=None, size=(1280, 720)):
    """Generate a clean slide-style image with modern gradient and text layout."""
    w, h = size

    # --- Create gradient background ---
    base_color = (230, 235, 255)
    accent_color = (190, 210, 255)
    img = Image.new('RGB', size, base_color)
    draw = ImageDraw.Draw(img)
    for y in range(h):
        ratio = y / h
        r = int(base_color[0] * (1 - ratio) + accent_color[0] * ratio)
        g = int(base_color[1] * (1 - ratio) + accent_color[1] * ratio)
        b = int(base_color[2] * (1 - ratio) + accent_color[2] * ratio)
        draw.line([(0, y), (w, y)], fill=(r, g, b))

    # --- Load fonts (fallback safe) ---
    try:
        font_title = ImageFont.truetype("arialbd.ttf", 64)
        font_bullet = ImageFont.truetype("arial.ttf", 36)
    except:
        font_title = ImageFont.load_default()
        font_bullet = ImageFont.load_default()

    # --- Draw title (centered + wrapped) ---
    title = text.strip() if text else "Untitled Slide"
    wrapped_title = textwrap.fill(title, width=25)
    # Compute text size using multiline_textbbox (modern Pillow)
    bbox = draw.multiline_textbbox((0, 0), wrapped_title, font=font_title, align='center')
    title_w = bbox[2] - bbox[0]
    title_h = bbox[3] - bbox[1]
    title_x = (w - title_w) / 2
    title_y = 100
    draw.multiline_text((title_x, title_y), wrapped_title, font=font_title, fill=(20, 20, 60), align='center')


    # --- Draw bullets ---
    if bullets:
        y_start = title_y + title_h + 80
        for i, b in enumerate(bullets[:6]):
            wrapped = textwrap.fill(b, width=40)
            draw.text((180, y_start + i * 60), f"• {wrapped}", font=font_bullet, fill=(40, 40, 70))

    # --- Add subtle vignette for depth ---
    overlay = Image.new('L', size, 0)
    for y in range(h):
        for x in range(w):
            dist = ((x - w / 2) ** 2 + (y - h / 2) ** 2) ** 0.5
            overlay.putpixel((x, y), int(255 * (dist / (w / 1.2))))
    img.putalpha(overlay.point(lambda p: 255 - min(p, 120)))
    img = img.convert("RGB")

    img.save(out_path)
    return out_path


def build_pptx(slides, pptx_path, images_dir):
    prs = Presentation()
    # Use a blank slide layout
    blank = prs.slide_layouts[6]
    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        left = Inches(0.5)
        top = Inches(0.4)
        # Title
        tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(1))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = s['headline']
        p.font.size = Pt(28)
        p.font.bold = True
        # Image
        img_path = s.get('image')
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.6), width=Inches(4.5))
        # Bullets
        bx = slide.shapes.add_textbox(Inches(5.5), Inches(2.0), Inches(4.0), Inches(3.5))
        btf = bx.text_frame
        btf.word_wrap = True
        for b in s.get('bullets', []):
            p = btf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(18)
        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = s.get('speaker_notes','')
    prs.save(pptx_path)
    return pptx_path


def synthesize_speech_per_slide(slides, out_audio_dir):
    aud_paths = []
    if not TTS_AVAILABLE:
        print('pyttsx3 not available; skipping TTS synthesis.')
        return aud_paths
    engine = pyttsx3.init()
    # Optionally change voice/rate
    rate = engine.getProperty('rate')
    engine.setProperty('rate', int(rate * 0.9))
    for i, s in enumerate(slides):
        txt = s.get('speaker_notes','')
        if not txt:
            txt = s.get('headline','')
        outp = os.path.join(out_audio_dir, f'slide_{i:02d}.wav')
        engine.save_to_file(txt, outp)
    engine.runAndWait()
    # pyttsx3 sometimes buffers; wait for files
    for i in range(len(slides)):
        p = os.path.join(out_audio_dir, f'slide_{i:02d}.wav')
        if os.path.exists(p):
            aud_paths.append(p)
    return aud_paths


def generate_background_music(duration_s, out_path, sr=22050):
    # Simple ambient loop: layered sine waves + soft noise
    t = np.linspace(0, duration_s, int(sr*duration_s))
    s1 = 0.06 * np.sin(2*np.pi*220*t)  # A3-ish
    s2 = 0.04 * np.sin(2*np.pi*330*t)  # E4-ish
    noise = 0.01 * np.random.normal(size=t.shape)
    s = s1 + s2 + noise
    # exponential fade in/out
    fade = np.linspace(0,1,len(s))
    s = s * fade
    # normalize
    s = s / (np.max(np.abs(s)) + 1e-9)
    # write wav via scipy
    try:
        from scipy.io.wavfile import write
        write(out_path, sr, (s * 32767).astype(np.int16))
    except Exception as e:
        print('scipy not available to write background music:', e)
        # fallback: create silent file
        import wave
        with wave.open(out_path,'w') as wf:
            wf.setnchannels(1); wf.setsampwidth(2); wf.setframerate(sr)
            wf.writeframes(b'\x00\x00' * int(sr*duration_s))
    return out_path


def make_video_from_slides(slides, images_dir, audio_paths, out_video_path, per_slide_duration=6):
    from moviepy import ImageClip, AudioFileClip, CompositeAudioClip, concatenate_videoclips
    
    clips = []
    audio_segments = []
    
    for i, s in enumerate(slides):
        img = s.get('image')
        if not img or not os.path.exists(img):
            img = generate_simple_image(
                s.get('headline', f"Slide {i+1}"),
                os.path.join(images_dir, f'slide_{i:02d}.png')
            )

        # Make ImageClip with duration
        clip = ImageClip(img, duration=per_slide_duration)
        clips.append(clip)
        
        # Attach corresponding audio if available
        apath = audio_paths[i] if i < len(audio_paths) else None
        if apath and os.path.exists(apath):
            # Calculate the start time
            start_time = i * per_slide_duration
            # Store audio file path and timing
            audio_segments.append((apath, start_time))

    # Concatenate video clips
    video = concatenate_videoclips(clips, method='compose')

    # Build composite audio manually
    audio_clips_list = []
    
    # Add speech audio segments
    for apath, start_time in audio_segments:
        audio_clips_list.append((AudioFileClip(apath), start_time))
    
    # Generate and add background music at lower volume
    bg_path = os.path.join(os.path.dirname(out_video_path), 'background_loop.wav')
    generate_background_music(math.ceil(video.duration)+1, bg_path)
    
    # For background music, we need to reduce volume
    # Read the audio file and create a modified version
    import wave
    import struct
    
    # Create a low-volume version of the background music
    bg_low_path = bg_path.replace('.wav', '_low.wav')
    try:
        with wave.open(bg_path, 'rb') as infile:
            params = infile.getparams()
            frames = infile.readframes(params.nframes)
            
            # Convert to numpy array and reduce volume
            audio_data = np.frombuffer(frames, dtype=np.int16)
            audio_data = (audio_data * 0.08).astype(np.int16)
            
            # Write low volume version
            with wave.open(bg_low_path, 'wb') as outfile:
                outfile.setparams(params)
                outfile.writeframes(audio_data.tobytes())
        
        audio_clips_list.append((AudioFileClip(bg_low_path), 0))
    except Exception as e:
        print(f"Warning: Could not reduce background music volume: {e}")
        audio_clips_list.append((AudioFileClip(bg_path), 0))
    
    # Combine all audio using CompositeAudioClip
    # In MoviePy 2.x, we pass list of tuples (clip, start_time)
    if audio_clips_list:
        combined_audio = CompositeAudioClip([clip.with_start(start) for clip, start in audio_clips_list])
        video = video.with_audio(combined_audio)
    
    video.write_videofile(out_video_path, fps=24, codec='libx264', audio_codec='aac')
    return out_video_path

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', required=True, help='input PDF file')
    parser.add_argument('--n_slides', type=int, default=8, help='number of slides to produce')
    parser.add_argument('--out', default='output', help='output folder')
    args = parser.parse_args()

    pdf_path = args.input
    n_slides = args.n_slides
    out_base = args.out
    ensure_outdir(out_base)
    pdf_stem = Path(pdf_path).stem
    out_project = out_base   # all outputs directly in output/
    images_dir = os.path.join(out_base, 'images')
    audio_dir = os.path.join(out_base, 'audio')
    ensure_outdir(images_dir)
    ensure_outdir(audio_dir)


    print('Extracting text...')
    paras = extract_paragraphs(pdf_path)
    if not paras:
        print('No text found in PDF. Exiting.')
        sys.exit(1)
    print(f'Found {len(paras)} paragraphs.')

    # Try embeddings-based section selection
    if USE_SBERT and len(paras) >= n_slides:
        print('Using sentence-transformers + KMeans for section selection...')
        try:
            sections = top_n_sections_with_embeddings(paras, n_slides)
        except Exception as e:
            print('Embedding path failed, falling back to sumy or heuristics:', e)
            sections = None
    else:
        sections = None

    if sections is None:
        if SUMY_AVAILABLE:
            print('Using Sumy LexRank summarizer as fallback...')
            try:
                summary_sentences = top_n_sections_with_sumy(paras, n_slides)
                sections = summary_sentences
            except Exception as e:
                print('Sumy failed:', e)
                sections = paras[:n_slides]
        else:
            print('No advanced summarizer available; using top paragraphs as sections...')
            sections = paras[:n_slides]

    # Build slide content
    slides = []
    for i, sec in enumerate(sections[:n_slides]):
        headline, bullets, speaker = make_headline_and_bullets(sec)
        if not headline:
            headline = 'Slide %d' % (i+1)
        img_path = os.path.join(images_dir, f'slide_{i:02d}.png')
        generate_simple_image(headline, img_path)
        slides.append({
            'headline': headline,
            'bullets': bullets,
            'speaker_notes': speaker,
            'image': img_path,
        })

    # PPTX
    pptx_path = os.path.join(out_project, f'{pdf_stem}.pptx')
    print('Building PPTX...')
    build_pptx(slides, pptx_path, images_dir)
    print('PPTX saved to', pptx_path)

    # TTS
    print('Synthesizing speech (TTS) if available...')
    audio_paths = synthesize_speech_per_slide(slides, audio_dir)

    # Make video
    video_path = os.path.join(out_project, f'{pdf_stem}.mp4')
    print('Generating MP4 video...')
    try:
        make_video_from_slides(slides, images_dir, audio_paths, video_path, per_slide_duration=6)
        print('Video saved to', video_path)
    except Exception as e:
        print('Video generation failed:', e)

    print('All done. Outputs in', out_project)


if __name__ == '__main__':
    main()
